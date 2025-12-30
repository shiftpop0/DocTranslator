# translate/powerpoint.py
"""
PPT文档翻译处理器 (.pptx)

核心功能：
1. 识别并区分标题/副标题/正文/表格等元素
2. 保留图片、图表等非文本元素
3. 智能调整容器大小和字体以适应译文
4. 防止元素重叠和超出边界
5. 双语模式正确复制幻灯片

元素处理策略：
- 标题/副标题：不换行，优先扩展宽度或缩小字体
- 正文文本框：可换行，可扩展高度
- 表格：保持大小，缩小字体
- 图片/图表：完全保留不处理
"""

import datetime
import logging
import re
import copy
from typing import List, Dict, Any, Optional, Tuple, Set
from dataclasses import dataclass
from threading import Event
from enum import Enum
from pptx import Presentation
from pptx.util import Pt, Emu, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.shapes.base import BaseShape
from pptx.shapes.placeholder import PlaceholderPicture
from pptx.text.text import TextFrame, _Paragraph, _Run
from pptx.oxml.ns import qn
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from lxml import etree

from . import to_translate
from . import common

# ==================== 配置 ====================

MAX_CHUNK_SIZE = 2000
MIN_FONT_SIZE = Pt(10)  # 最小字体
MAX_WIDTH_EXPANSION = 1.3  # 最大宽度扩展比例
MAX_HEIGHT_EXPANSION = 1.5  # 最大高度扩展比例
FONT_SHRINK_STEP = 0.9  # 字体缩小步长
MIN_FONT_SCALE = 0.6  # 最小字体缩放比例


class ElementType(Enum):
    """元素类型"""
    TITLE = "title"
    SUBTITLE = "subtitle"
    BODY = "body"
    TABLE_CELL = "table_cell"
    TEXT_BOX = "text_box"
    OTHER = "other"


@dataclass
class RunStyle:
    """Run样式"""
    font_name: Optional[str] = None
    font_name_ea: Optional[str] = None
    font_size: Optional[int] = None
    bold: Optional[bool] = None
    italic: Optional[bool] = None
    underline: Optional[bool] = None
    color_rgb: Optional[RGBColor] = None


@dataclass
class ShapeGeometry:
    """形状几何信息"""
    left: int = 0
    top: int = 0
    width: int = 0
    height: int = 0


@dataclass
class TextBlock:
    """文本块"""
    uid: str
    slide_index: int
    shape_id: int
    element_type: ElementType

    # 位置信息
    location_type: str = "textframe"  # textframe, table_cell
    paragraph_index: int = 0
    cell_row: int = -1
    cell_col: int = -1

    # 文本内容
    original_text: str = ""
    translated_text: str = ""

    # 样式和几何
    run_style: Optional[RunStyle] = None
    geometry: Optional[ShapeGeometry] = None

    # 状态
    complete: bool = False
    skip: bool = False
    count: int = 0

    # 子块信息
    is_sub: bool = False
    sub_index: int = 0
    sub_total: int = 1
    parent_uid: str = ""


# ==================== 入口函数 ====================

def start(trans: Dict[str, Any]) -> bool:
    """PPT翻译入口"""
    translate_id = trans['id']
    start_time = datetime.datetime.now()

    trans_type = trans.get('type', 'trans_only_inherit')
    only_translation = 'only' in trans_type
    is_bilingual = 'both' in trans_type
    target_lang = trans.get('lang', '英语')

    logging.info(
        f"[任务{translate_id}] PPT翻译模式: only={only_translation}, bilingual={is_bilingual}")

    # 打开文件
    try:
        prs = Presentation(trans['file_path'])
    except Exception as e:
        logging.error(f"[任务{translate_id}] 打开PPT失败: {e}")
        to_translate.error(translate_id, f"打开PPT失败: {str(e)}")
        return False

    # 获取幻灯片尺寸（用于边界检测）
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # 提取文本块
    try:
        all_blocks = _extract_all_blocks(prs)
    except Exception as e:
        logging.error(f"[任务{translate_id}] 提取文本失败: {e}")
        to_translate.error(translate_id, f"提取文本失败: {str(e)}")
        return False

    to_translate_blocks = [b for b in all_blocks if not b.skip]

    if not to_translate_blocks:
        logging.info(f"[任务{translate_id}] 没有需要翻译的文本")
        prs.save(trans['target_file'])
        to_translate.complete(trans, 0, "0秒")
        return True

    logging.info(f"[任务{translate_id}] 共{len(all_blocks)}块，{len(to_translate_blocks)}块需翻译")

    # 执行翻译
    texts = _blocks_to_api_format(to_translate_blocks)
    event = Event()
    success = to_translate.translate_batch(trans, texts, event)

    if not success:
        return False

    _sync_translation_results(to_translate_blocks, texts)

    # 应用翻译
    try:
        if is_bilingual:
            text_count = _apply_bilingual_mode(prs, all_blocks, target_lang,
                                               slide_width, slide_height)
        else:
            text_count = _apply_translation_mode(prs, all_blocks, target_lang,
                                                 slide_width, slide_height)

        prs.save(trans['target_file'])
    except Exception as e:
        logging.error(f"[任务{translate_id}] 保存失败: {e}")
        to_translate.error(translate_id, f"保存失败: {str(e)}")
        return False

    end_time = datetime.datetime.now()
    spend_time = common.display_spend(start_time, end_time)
    to_translate.complete(trans, text_count, spend_time)
    return True


# ==================== 文本提取 ====================

def _extract_all_blocks(prs: Presentation) -> List[TextBlock]:
    """提取所有文本块"""
    blocks = []
    uid_counter = [0]

    def next_uid() -> str:
        uid_counter[0] += 1
        return f"blk_{uid_counter[0]}"

    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            shape_blocks = _extract_shape_blocks(shape, slide_idx, next_uid)
            blocks.extend(shape_blocks)

    return blocks


def _extract_shape_blocks(shape: BaseShape, slide_idx: int, next_uid) -> List[TextBlock]:
    """提取形状中的文本块"""
    blocks = []
    shape_id = shape.shape_id

    # 1. 跳过图片
    if _is_picture(shape):
        logging.debug(f"跳过图片: shape_id={shape_id}")
        return blocks

    # 2. 跳过图表
    if _is_chart(shape):
        logging.debug(f"跳过图表: shape_id={shape_id}")
        return blocks

    # 3. 跳过媒体/OLE对象
    try:
        skip_types = {
            MSO_SHAPE_TYPE.MEDIA,
            MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT,
            MSO_SHAPE_TYPE.LINKED_OLE_OBJECT,
            MSO_SHAPE_TYPE.IGX_GRAPHIC,
        }
        if shape.shape_type in skip_types:
            return blocks
    except:
        pass

    # 4. 组合形状递归处理
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        try:
            for sub_shape in shape.shapes:
                sub_blocks = _extract_shape_blocks(sub_shape, slide_idx, next_uid)
                blocks.extend(sub_blocks)
        except:
            pass
        return blocks

    # 5. 获取形状几何信息
    geometry = _get_shape_geometry(shape)

    # 6. 识别元素类型
    element_type = _identify_element_type(shape)

    # 7. 表格
    if shape.has_table:
        blocks.extend(_extract_table_blocks(shape.table, slide_idx, shape_id,
                                            geometry, next_uid))
        return blocks

    # 8. 文本框
    if shape.has_text_frame:
        blocks.extend(_extract_textframe_blocks(shape.text_frame, slide_idx, shape_id,
                                                element_type, geometry, next_uid))

    return blocks


def _is_picture(shape: BaseShape) -> bool:
    """判断是否是图片"""
    # 检查shape类型
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return True
    except:
        pass

    # 检查是否是图片占位符
    try:
        if isinstance(shape, PlaceholderPicture):
            return True
    except:
        pass

    # 检查XML中是否包含图片
    try:
        elem = shape._element
        # 检查blip（图片引用）
        blips = elem.findall('.//' + qn('a:blip'))
        if blips:
            # 如果有blip但没有文本框，是纯图片
            if not shape.has_text_frame:
                return True
            # 如果有文本框但为空，也当作图片
            if shape.has_text_frame and not shape.text_frame.text.strip():
                return True
    except:
        pass

    return False


def _is_chart(shape: BaseShape) -> bool:
    """判断是否是图表"""
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.CHART:
            return True
    except:
        pass

    # 检查XML
    try:
        elem = shape._element
        charts = elem.findall('.//' + qn('c:chart'))
        if charts:
            return True
    except:
        pass

    try:
        if hasattr(shape, 'chart'):
            return True
    except:
        pass

    return False


def _get_shape_geometry(shape: BaseShape) -> ShapeGeometry:
    """获取形状几何信息"""
    try:
        return ShapeGeometry(
            left=shape.left or 0,
            top=shape.top or 0,
            width=shape.width or 0,
            height=shape.height or 0
        )
    except:
        return ShapeGeometry()


def _identify_element_type(shape: BaseShape) -> ElementType:
    """识别元素类型"""
    # 检查是否是占位符
    try:
        if shape.is_placeholder:
            ph_type = shape.placeholder_format.type

            # 标题类型
            if ph_type in [PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE]:
                return ElementType.TITLE

            # 副标题类型
            if ph_type in [PP_PLACEHOLDER.SUBTITLE]:
                return ElementType.SUBTITLE

            # 正文类型
            if ph_type in [PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT]:
                return ElementType.BODY
    except:
        pass

    # 非占位符的文本框
    if shape.has_text_frame:
        return ElementType.TEXT_BOX

    return ElementType.OTHER


def _extract_textframe_blocks(text_frame: TextFrame, slide_idx: int, shape_id: int,
                              element_type: ElementType, geometry: ShapeGeometry,
                              next_uid) -> List[TextBlock]:
    """提取文本框中的段落"""
    blocks = []

    for para_idx, paragraph in enumerate(text_frame.paragraphs):
        text = _get_paragraph_text(paragraph)

        if not text or not text.strip():
            continue

        run_style = _extract_first_run_style(paragraph)
        should_skip = not _should_translate(text)

        if should_skip:
            blocks.append(TextBlock(
                uid=next_uid(),
                slide_index=slide_idx,
                shape_id=shape_id,
                element_type=element_type,
                location_type='textframe',
                paragraph_index=para_idx,
                original_text=text,
                run_style=run_style,
                geometry=geometry,
                skip=True
            ))
            continue

        # 分块
        if len(text) <= MAX_CHUNK_SIZE:
            blocks.append(TextBlock(
                uid=next_uid(),
                slide_index=slide_idx,
                shape_id=shape_id,
                element_type=element_type,
                location_type='textframe',
                paragraph_index=para_idx,
                original_text=text,
                run_style=run_style,
                geometry=geometry
            ))
        else:
            sub_texts = _split_by_sentences(text, MAX_CHUNK_SIZE)
            parent_uid = next_uid()
            for i, sub_text in enumerate(sub_texts):
                blocks.append(TextBlock(
                    uid=next_uid(),
                    slide_index=slide_idx,
                    shape_id=shape_id,
                    element_type=element_type,
                    location_type='textframe',
                    paragraph_index=para_idx,
                    original_text=sub_text,
                    run_style=run_style,
                    geometry=geometry,
                    is_sub=True,
                    sub_index=i,
                    sub_total=len(sub_texts),
                    parent_uid=parent_uid
                ))

    return blocks


def _extract_table_blocks(table, slide_idx: int, shape_id: int,
                          geometry: ShapeGeometry, next_uid) -> List[TextBlock]:
    """提取表格文本"""
    blocks = []
    processed_cells: Set[int] = set()

    for row_idx, row in enumerate(table.rows):
        for col_idx, cell in enumerate(row.cells):
            cell_id = id(cell._tc)
            if cell_id in processed_cells:
                continue
            processed_cells.add(cell_id)

            cell_text = _get_cell_text(cell)

            if not cell_text or not cell_text.strip():
                continue

            run_style = None
            if cell.text_frame.paragraphs:
                run_style = _extract_first_run_style(cell.text_frame.paragraphs[0])

            should_skip = not _should_translate(cell_text)

            if should_skip:
                blocks.append(TextBlock(
                    uid=next_uid(),
                    slide_index=slide_idx,
                    shape_id=shape_id,
                    element_type=ElementType.TABLE_CELL,
                    location_type='table_cell',
                    cell_row=row_idx,
                    cell_col=col_idx,
                    original_text=cell_text,
                    run_style=run_style,
                    geometry=geometry,
                    skip=True
                ))
                continue

            if len(cell_text) <= MAX_CHUNK_SIZE:
                blocks.append(TextBlock(
                    uid=next_uid(),
                    slide_index=slide_idx,
                    shape_id=shape_id,
                    element_type=ElementType.TABLE_CELL,
                    location_type='table_cell',
                    cell_row=row_idx,
                    cell_col=col_idx,
                    original_text=cell_text,
                    run_style=run_style,
                    geometry=geometry
                ))
            else:
                sub_texts = _split_by_sentences(cell_text, MAX_CHUNK_SIZE)
                parent_uid = next_uid()
                for i, sub_text in enumerate(sub_texts):
                    blocks.append(TextBlock(
                        uid=next_uid(),
                        slide_index=slide_idx,
                        shape_id=shape_id,
                        element_type=ElementType.TABLE_CELL,
                        location_type='table_cell',
                        cell_row=row_idx,
                        cell_col=col_idx,
                        original_text=sub_text,
                        run_style=run_style,
                        geometry=geometry,
                        is_sub=True,
                        sub_index=i,
                        sub_total=len(sub_texts),
                        parent_uid=parent_uid
                    ))

    return blocks


def _get_paragraph_text(paragraph: _Paragraph) -> str:
    """获取段落文本"""
    parts = []
    for run in paragraph.runs:
        if run.text:
            parts.append(run.text)
    return ''.join(parts)


def _get_cell_text(cell) -> str:
    """获取单元格文本"""
    parts = []
    for para in cell.text_frame.paragraphs:
        para_text = _get_paragraph_text(para)
        if para_text:
            parts.append(para_text)
    return '\n'.join(parts)


def _extract_first_run_style(paragraph: _Paragraph) -> Optional[RunStyle]:
    """提取段落第一个有效run的样式"""
    for run in paragraph.runs:
        if run.text and run.text.strip():
            return _extract_run_style(run)
    return None


def _extract_run_style(run: _Run) -> RunStyle:
    """提取run样式"""
    style = RunStyle()

    try:
        font = run.font
        style.font_name = font.name

        # 东亚字体
        try:
            rPr = run._r.get_or_add_rPr()
            ea = rPr.find(qn('a:ea'))
            if ea is not None:
                style.font_name_ea = ea.get('typeface')
        except:
            pass

        if font.size:
            style.font_size = font.size

        style.bold = font.bold
        style.italic = font.italic
        style.underline = font.underline

        try:
            if font.color and font.color.rgb:
                style.color_rgb = font.color.rgb
        except:
            pass

    except:
        pass

    return style


def _should_translate(text: str) -> bool:
    """判断是否需要翻译"""
    text = text.strip()

    if not text:
        return False

    if len(text) < 2:
        return False

    if common.is_all_punc(text):
        return False

    # 纯数字/符号
    if re.match(r'^[\d\s\.\-\+\*\/\=\%\(\)\[\]\{\}\#\@\&\|\\\:\;\,\<\>\$\€\¥\£]+$', text):
        return False

    # 页码
    if re.match(r'^(第?\s*\d+\s*页?|Page\s*\d+|\d+\s*/\s*\d+)$', text, re.IGNORECASE):
        return False

    # 日期时间
    if re.match(r'^\d{4}[-/]\d{1,2}[-/]\d{1,2}', text):
        return False

    # URL
    if re.match(r'^https?://', text):
        return False

    # 邮箱
    if re.match(r'^[\w\.-]+@[\w\.-]+\.\w+$', text):
        return False

    return True


def _split_by_sentences(text: str, max_size: int) -> List[str]:
    """按句子边界切分"""
    endings = r'([.!?。！？；;]\s*)'
    parts = re.split(endings, text)

    sentences = []
    i = 0
    while i < len(parts):
        s = parts[i]
        if i + 1 < len(parts) and re.match(endings, parts[i + 1]):
            s += parts[i + 1]
            i += 2
        else:
            i += 1
        if s.strip():
            sentences.append(s)

    chunks = []
    current = ""

    for s in sentences:
        if len(s) > max_size:
            if current:
                chunks.append(current)
                current = ""
            for j in range(0, len(s), max_size):
                chunks.append(s[j:j + max_size])
        elif len(current) + len(s) <= max_size:
            current += s
        else:
            if current:
                chunks.append(current)
            current = s

    if current:
        chunks.append(current)

    return chunks if chunks else [text]


# ==================== 翻译接口 ====================

def _blocks_to_api_format(blocks: List[TextBlock]) -> List[Dict]:
    """转换为API格式"""
    return [{
        'text': b.original_text,
        'original': b.original_text,
        'complete': False,
        'count': 0,
        '_uid': b.uid
    } for b in blocks]


def _sync_translation_results(blocks: List[TextBlock], texts: List[Dict]):
    """同步翻译结果"""
    uid_map = {b.uid: b for b in blocks}
    for t in texts:
        uid = t.get('_uid')
        if uid and uid in uid_map:
            block = uid_map[uid]
            block.translated_text = t.get('text', block.original_text)
            block.count = t.get('count', 0)
            block.complete = True


# ==================== 仅译文模式 ====================

def _apply_translation_mode(prs: Presentation, blocks: List[TextBlock],
                            target_lang: str, slide_width: int, slide_height: int) -> int:
    """应用仅译文模式"""
    text_count = 0

    # 按幻灯片分组
    slide_blocks = _group_by_slide(blocks)

    for slide_idx, slide in enumerate(prs.slides):
        if slide_idx not in slide_blocks:
            continue

        sblocks = slide_blocks[slide_idx]
        shape_map = _build_shape_map(slide)

        # 获取所有形状的几何信息（用于碰撞检测）
        all_geometries = _get_all_shape_geometries(slide)

        # 按shape分组
        shape_blocks = _group_by_shape(sblocks)

        for shape_id, blocks_list in shape_blocks.items():
            if shape_id not in shape_map:
                continue

            shape = shape_map[shape_id]
            count = _apply_to_shape(shape, blocks_list, target_lang,
                                    slide_width, slide_height, all_geometries)
            text_count += count

    return text_count


def _group_by_slide(blocks: List[TextBlock]) -> Dict[int, List[TextBlock]]:
    """按幻灯片分组"""
    result = {}
    for b in blocks:
        if b.slide_index not in result:
            result[b.slide_index] = []
        result[b.slide_index].append(b)
    return result


def _group_by_shape(blocks: List[TextBlock]) -> Dict[int, List[TextBlock]]:
    """按shape分组"""
    result = {}
    for b in blocks:
        if b.shape_id not in result:
            result[b.shape_id] = []
        result[b.shape_id].append(b)
    return result


def _build_shape_map(slide) -> Dict[int, BaseShape]:
    """构建shape_id到shape的映射"""
    shape_map = {}

    def add_shape(shape):
        shape_map[shape.shape_id] = shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            try:
                for sub_shape in shape.shapes:
                    add_shape(sub_shape)
            except:
                pass

    for shape in slide.shapes:
        add_shape(shape)

    return shape_map


def _get_all_shape_geometries(slide) -> List[ShapeGeometry]:
    """获取幻灯片上所有形状的几何信息"""
    geometries = []

    def collect_geometry(shape):
        try:
            geo = _get_shape_geometry(shape)
            if geo.width > 0 and geo.height > 0:
                geometries.append(geo)
        except:
            pass

        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            try:
                for sub_shape in shape.shapes:
                    collect_geometry(sub_shape)
            except:
                pass

    for shape in slide.shapes:
        collect_geometry(shape)

    return geometries


def _apply_to_shape(shape: BaseShape, blocks: List[TextBlock], target_lang: str,
                    slide_width: int, slide_height: int,
                    all_geometries: List[ShapeGeometry]) -> int:
    """应用翻译到形状"""
    if shape.has_table:
        return _apply_to_table(shape.table, blocks, target_lang)
    elif shape.has_text_frame:
        return _apply_to_textframe(shape, blocks, target_lang,
                                   slide_width, slide_height, all_geometries)
    return 0


def _apply_to_textframe(shape: BaseShape, blocks: List[TextBlock], target_lang: str,
                        slide_width: int, slide_height: int,
                        all_geometries: List[ShapeGeometry]) -> int:
    """应用翻译到文本框"""
    text_count = 0
    text_frame = shape.text_frame

    # 按段落分组
    para_blocks = {}
    for b in blocks:
        if b.location_type != 'textframe':
            continue
        if b.paragraph_index not in para_blocks:
            para_blocks[b.paragraph_index] = []
        para_blocks[b.paragraph_index].append(b)

    # 处理每个段落
    for para_idx, paragraph in enumerate(text_frame.paragraphs):
        if para_idx not in para_blocks:
            continue

        pblocks = sorted(para_blocks[para_idx], key=lambda x: x.sub_index)

        if pblocks[0].skip:
            continue

        # 合并子块
        if any(b.is_sub for b in pblocks):
            translated = ''.join(b.translated_text or b.original_text for b in pblocks)
            original_text = ''.join(b.original_text for b in pblocks)
        else:
            translated = pblocks[0].translated_text or pblocks[0].original_text
            original_text = pblocks[0].original_text

        text_count += sum(b.count for b in pblocks)
        element_type = pblocks[0].element_type
        run_style = pblocks[0].run_style

        # 根据元素类型决定处理方式
        _replace_paragraph_text_smart(
            paragraph=paragraph,
            shape=shape,
            new_text=translated,
            original_text=original_text,
            element_type=element_type,
            run_style=run_style,
            target_lang=target_lang,
            slide_width=slide_width,
            slide_height=slide_height,
            all_geometries=all_geometries
        )

    return text_count


def _apply_to_table(table, blocks: List[TextBlock], target_lang: str) -> int:
    """应用翻译到表格"""
    text_count = 0

    # 按单元格分组
    cell_blocks = {}
    for b in blocks:
        if b.location_type != 'table_cell':
            continue
        key = (b.cell_row, b.cell_col)
        if key not in cell_blocks:
            cell_blocks[key] = []
        cell_blocks[key].append(b)

    for (row_idx, col_idx), cblocks in cell_blocks.items():
        try:
            cell = table.cell(row_idx, col_idx)
        except:
            continue

        cblocks = sorted(cblocks, key=lambda x: x.sub_index)

        if cblocks[0].skip:
            continue

        if any(b.is_sub for b in cblocks):
            translated = ''.join(b.translated_text or b.original_text for b in cblocks)
            original_text = ''.join(b.original_text for b in cblocks)
        else:
            translated = cblocks[0].translated_text or cblocks[0].original_text
            original_text = cblocks[0].original_text

        text_count += sum(b.count for b in cblocks)

        _replace_cell_text_smart(cell, translated, original_text,
                                 cblocks[0].run_style, target_lang)

    return text_count


def _replace_paragraph_text_smart(paragraph: _Paragraph, shape: BaseShape,
                                  new_text: str, original_text: str,
                                  element_type: ElementType,
                                  run_style: Optional[RunStyle],
                                  target_lang: str,
                                  slide_width: int, slide_height: int,
                                  all_geometries: List[ShapeGeometry]):
    """
    智能替换段落文本
    根据元素类型和文本长度变化调整容器和字体
    """
    runs = list(paragraph.runs)

    if not runs:
        run = paragraph.add_run()
        run.text = new_text
        _apply_run_style(run, run_style, target_lang)
        return

    # 计算长度变化
    original_len = len(original_text)
    new_len = len(new_text)
    length_ratio = new_len / max(original_len, 1)

    # 替换文本
    first_run = runs[0]
    first_run.text = new_text
    for run in runs[1:]:
        run.text = ""

    # 确保字体兼容
    _ensure_font_compatibility(first_run, target_lang, run_style)

    # 如果文本变长，需要调整
    if length_ratio > 1.1:
        _adjust_for_longer_text(
            paragraph=paragraph,
            shape=shape,
            first_run=first_run,
            length_ratio=length_ratio,
            element_type=element_type,
            run_style=run_style,
            slide_width=slide_width,
            slide_height=slide_height,
            all_geometries=all_geometries
        )


def _adjust_for_longer_text(paragraph: _Paragraph, shape: BaseShape, first_run: _Run,
                            length_ratio: float, element_type: ElementType,
                            run_style: Optional[RunStyle],
                            slide_width: int, slide_height: int,
                            all_geometries: List[ShapeGeometry]):
    """
    调整以适应更长的文本
    策略：
    1. 标题/副标题：优先缩小字体，其次扩展宽度，避免换行
    2. 正文/文本框：优先扩展高度，其次缩小字体，允许换行
    3. 所有调整都要检查边界和碰撞
    """

    if element_type in [ElementType.TITLE, ElementType.SUBTITLE]:
        # 标题类：优先缩小字体
        _adjust_title_element(shape, first_run, length_ratio, run_style,
                              slide_width, slide_height, all_geometries)
    else:
        # 正文类：优先扩展容器
        _adjust_body_element(shape, first_run, length_ratio, run_style,
                             slide_width, slide_height, all_geometries)


def _adjust_title_element(shape: BaseShape, run: _Run, length_ratio: float,
                          run_style: Optional[RunStyle],
                          slide_width: int, slide_height: int,
                          all_geometries: List[ShapeGeometry]):
    """
    调整标题元素
    策略：缩小字体 > 扩展宽度 > 移动位置
    """
    # 1. 首先尝试缩小字体
    font_scale = _calculate_font_scale_for_title(length_ratio)
    if font_scale < 1.0:
        _scale_run_font(run, font_scale, run_style)

    # 2. 如果还是太长，尝试扩展宽度
    if length_ratio > 1.5 and font_scale >= MIN_FONT_SCALE:
        try:
            current_width = shape.width
            current_left = shape.left

            # 计算需要的新宽度
            needed_width = int(current_width * min(length_ratio * 0.8, MAX_WIDTH_EXPANSION))

            # 检查右边界
            max_width = slide_width - current_left - Emu(Inches(0.2))  # 留边距
            new_width = min(needed_width, max_width)

            # 检查是否会与其他元素重叠
            new_geo = ShapeGeometry(
                left=current_left,
                top=shape.top,
                width=new_width,
                height=shape.height
            )

            if not _would_overlap(new_geo, all_geometries, shape.shape_id, shape):
                shape.width = new_width
            else:
                # 如果扩展会重叠，进一步缩小字体
                _scale_run_font(run, 0.85, run_style)

        except Exception as e:
            logging.debug(f"调整标题宽度失败: {e}")


def _adjust_body_element(shape: BaseShape, run: _Run, length_ratio: float,
                         run_style: Optional[RunStyle],
                         slide_width: int, slide_height: int,
                         all_geometries: List[ShapeGeometry]):
    """
    调整正文元素
    策略：扩展高度 > 缩小字体
    """
    # 1. 首先尝试扩展高度（允许换行）
    if length_ratio > 1.2:
        try:
            current_height = shape.height

            # 计算需要的新高度
            needed_height = int(current_height * min(length_ratio * 0.9, MAX_HEIGHT_EXPANSION))

            # 检查下边界
            max_height = slide_height - shape.top - Emu(Inches(0.2))
            new_height = min(needed_height, max_height)

            # 检查是否会与其他元素重叠
            new_geo = ShapeGeometry(
                left=shape.left,
                top=shape.top,
                width=shape.width,
                height=new_height
            )

            if not _would_overlap(new_geo, all_geometries, shape.shape_id, shape):
                shape.height = new_height
            else:
                # 如果扩展会重叠，缩小字体
                font_scale = _calculate_font_scale_for_body(length_ratio)
                _scale_run_font(run, font_scale, run_style)

        except Exception as e:
            logging.debug(f"调整正文高度失败: {e}")
            # 失败时缩小字体
            font_scale = _calculate_font_scale_for_body(length_ratio)
            _scale_run_font(run, font_scale, run_style)

    # 2. 如果比例过大，还需要缩小字体
    if length_ratio > 1.8:
        font_scale = _calculate_font_scale_for_body(length_ratio)
        _scale_run_font(run, font_scale, run_style)


def _calculate_font_scale_for_title(length_ratio: float) -> float:
    """计算标题的字体缩放比例"""
    if length_ratio <= 1.2:
        return 1.0
    elif length_ratio <= 1.5:
        return 0.90
    elif length_ratio <= 2.0:
        return 0.80
    elif length_ratio <= 2.5:
        return 0.70
    else:
        return MIN_FONT_SCALE


def _calculate_font_scale_for_body(length_ratio: float) -> float:
    """计算正文的字体缩放比例"""
    if length_ratio <= 1.5:
        return 1.0
    elif length_ratio <= 2.0:
        return 0.95
    elif length_ratio <= 2.5:
        return 0.90
    elif length_ratio <= 3.0:
        return 0.85
    else:
        return 0.75


def _scale_run_font(run: _Run, scale: float, run_style: Optional[RunStyle]):
    """缩放run的字体大小"""
    try:
        # 获取当前字体大小
        current_size = run.font.size
        if current_size is None and run_style and run_style.font_size:
            current_size = run_style.font_size
        if current_size is None:
            current_size = Pt(18)  # 默认大小

        # 计算新大小
        new_size = int(current_size * scale)
        if new_size < MIN_FONT_SIZE:
            new_size = MIN_FONT_SIZE

        run.font.size = new_size
    except Exception as e:
        logging.debug(f"缩放字体失败: {e}")


def _would_overlap(new_geo: ShapeGeometry, all_geometries: List[ShapeGeometry],
                   exclude_shape_id: int, shape: BaseShape) -> bool:
    """检查新几何是否会与其他形状重叠"""
    # 获取当前形状的几何
    current_geo = _get_shape_geometry(shape)

    for geo in all_geometries:
        # 跳过自己
        if (geo.left == current_geo.left and geo.top == current_geo.top and
                geo.width == current_geo.width and geo.height == current_geo.height):
            continue

        # 检查是否重叠
        if _geometries_overlap(new_geo, geo):
            return True

    return False


def _geometries_overlap(geo1: ShapeGeometry, geo2: ShapeGeometry) -> bool:
    """检查两个几何是否重叠"""
    # AABB碰撞检测
    left1, top1, right1, bottom1 = geo1.left, geo1.top, geo1.left + geo1.width, geo1.top + geo1.height
    left2, top2, right2, bottom2 = geo2.left, geo2.top, geo2.left + geo2.width, geo2.top + geo2.height

    # 如果不重叠，返回False
    if right1 <= left2 or right2 <= left1:
        return False
    if bottom1 <= top2 or bottom2 <= top1:
        return False

    return True


def _replace_cell_text_smart(cell, new_text: str, original_text: str,
                             run_style: Optional[RunStyle], target_lang: str):
    """智能替换单元格文本"""
    paragraphs = cell.text_frame.paragraphs
    lines = new_text.split('\n')

    # 计算长度变化
    length_ratio = len(new_text) / max(len(original_text), 1)
    font_scale = 1.0
    if length_ratio > 1.3:
        font_scale = max(0.7, 1.0 / (length_ratio * 0.8))

    for i, line in enumerate(lines):
        if i < len(paragraphs):
            para = paragraphs[i]
            runs = list(para.runs)

            if runs:
                runs[0].text = line
                for run in runs[1:]:
                    run.text = ""

                if font_scale < 1.0:
                    _scale_run_font(runs[0], font_scale, run_style)
                _ensure_font_compatibility(runs[0], target_lang, run_style)
            else:
                run = para.add_run()
                run.text = line
                if font_scale < 1.0:
                    _scale_run_font(run, font_scale, run_style)
                _apply_run_style(run, run_style, target_lang)
        elif paragraphs:
            # 追加到最后一个段落
            last_para = paragraphs[-1]
            if last_para.runs:
                last_para.runs[-1].text += '\n' + line


def _apply_run_style(run: _Run, style: Optional[RunStyle], target_lang: str):
    """应用run样式"""
    if not style:
        _set_default_font(run, target_lang)
        return

    try:
        font = run.font

        if style.bold is not None:
            font.bold = style.bold
        if style.italic is not None:
            font.italic = style.italic
        if style.underline is not None:
            font.underline = style.underline
        if style.font_size:
            font.size = style.font_size
        if style.color_rgb:
            font.color.rgb = style.color_rgb

        _set_font_names(run, style, target_lang)

    except:
        pass


def _ensure_font_compatibility(run: _Run, target_lang: str, style: Optional[RunStyle]):
    """确保字体兼容目标语言"""
    try:
        current_font = run.font.name

        if target_lang in ['中文', '日语', '韩语']:
            if not _is_cjk_font(current_font):
                _set_font_names(run, style, target_lang)
    except:
        pass


def _set_font_names(run: _Run, style: Optional[RunStyle], target_lang: str):
    """设置字体名称"""
    try:
        # 西文字体
        latin_font = 'Arial'
        if style and style.font_name:
            latin_font = style.font_name

        # 东亚字体
        ea_font = None
        if style and style.font_name_ea:
            ea_font = style.font_name_ea

        # 根据目标语言设置东亚字体
        if target_lang in ['中文', '日语', '韩语']:
            if not ea_font or not _is_cjk_font(ea_font):
                if target_lang == '中文':
                    ea_font = 'Microsoft YaHei'
                elif target_lang == '日语':
                    ea_font = 'Yu Gothic'
                elif target_lang == '韩语':
                    ea_font = 'Malgun Gothic'

        run.font.name = latin_font

        # 设置东亚字体
        if ea_font:
            try:
                rPr = run._r.get_or_add_rPr()

                # 设置latin字体
                latin = rPr.find(qn('a:latin'))
                if latin is None:
                    latin = etree.SubElement(rPr, qn('a:latin'))
                latin.set('typeface', latin_font)

                # 设置ea字体
                ea = rPr.find(qn('a:ea'))
                if ea is None:
                    ea = etree.SubElement(rPr, qn('a:ea'))
                ea.set('typeface', ea_font)
            except:
                pass

    except:
        pass


def _set_default_font(run: _Run, target_lang: str):
    """设置默认字体"""
    try:
        if target_lang in ['中文', '日语', '韩语']:
            run.font.name = 'Microsoft YaHei'
        else:
            run.font.name = 'Arial'
    except:
        pass


def _is_cjk_font(font_name: str) -> bool:
    """检查是否是CJK字体"""
    if not font_name:
        return False

    cjk_keywords = [
        'yahei', '雅黑', 'simsun', '宋体', 'simhei', '黑体',
        'kaiti', '楷体', 'fangsong', '仿宋',
        'gothic', 'mincho', 'meiryo', 'hiragino',
        'malgun', 'batang', 'gulim', 'dotum',
        'noto sans cjk', 'noto serif cjk', 'source han',
        'pingfang', '苹方', 'heiti', 'songti',
        'microsoft jhenghei', '微軟正黑'
    ]

    font_lower = font_name.lower()
    return any(kw in font_lower for kw in cjk_keywords)


# ==================== 双语模式 ====================

def _apply_bilingual_mode(prs: Presentation, blocks: List[TextBlock],
                          target_lang: str, slide_width: int, slide_height: int) -> int:
    """
    双语模式：每个原文幻灯片后插入译文幻灯片
    """
    text_count = 0

    # 按幻灯片分组
    slide_blocks = _group_by_slide(blocks)

    original_count = len(prs.slides)

    # 从后往前处理
    for slide_idx in range(original_count - 1, -1, -1):
        try:
            original_slide = prs.slides[slide_idx]

            # 复制幻灯片（包含图片等所有元素）
            new_slide = _duplicate_slide(prs, original_slide)

            if new_slide is None:
                logging.warning(f"复制幻灯片 {slide_idx} 失败")
                continue

            # 移动到正确位置
            _move_slide(prs, len(prs.slides) - 1, slide_idx + 1)

            # 获取新幻灯片并应用翻译
            translated_slide = prs.slides[slide_idx + 1]

            if slide_idx in slide_blocks:
                sblocks = slide_blocks[slide_idx]

                # 建立位置映射
                shape_mapping = _build_shape_mapping_by_position(original_slide, translated_slide)

                # 统计
                for b in sblocks:
                    if not b.skip:
                        text_count += b.count

                # 应用翻译
                _apply_blocks_with_mapping(translated_slide, sblocks, shape_mapping,
                                           target_lang, slide_width, slide_height)

        except Exception as e:
            logging.error(f"双语模式处理幻灯片 {slide_idx} 失败: {e}")
            import traceback
            traceback.print_exc()

    return text_count


def _duplicate_slide(prs: Presentation, source_slide) -> Optional[object]:
    """
    完整复制幻灯片，包括所有元素和图片关系

    关键：正确处理图片的relationship引用
    """
    try:
        # 使用源幻灯片的布局创建新幻灯片
        new_slide = prs.slides.add_slide(source_slide.slide_layout)

        # 获取源幻灯片和新幻灯片的part（用于处理关系）
        source_part = source_slide.part
        new_part = new_slide.part

        # 清除新幻灯片上由布局自动生成的形状
        spTree = new_slide.shapes._spTree
        shapes_to_remove = []

        for child in spTree:
            tag = child.tag
            # 保留 nvGrpSpPr 和 grpSpPr（组属性），删除其他形状
            if tag.endswith('}nvGrpSpPr') or tag.endswith('}grpSpPr'):
                continue
            if (tag.endswith('}sp') or tag.endswith('}pic') or
                    tag.endswith('}graphicFrame') or tag.endswith('}grpSp') or
                    tag.endswith('}cxnSp')):
                shapes_to_remove.append(child)

        for elem in shapes_to_remove:
            spTree.remove(elem)

        # 复制源幻灯片的所有形状，并正确处理关系
        for shape in source_slide.shapes:
            _copy_shape_with_relationships(shape, spTree, source_part, new_part)

        # 复制背景
        _copy_slide_background(source_slide, new_slide)

        return new_slide

    except Exception as e:
        logging.error(f"复制幻灯片失败: {e}")
        import traceback
        traceback.print_exc()
        return None


def _copy_shape_with_relationships(shape: BaseShape, target_spTree,
                                   source_part, target_part):
    """
    复制形状及其关联的资源（图片等）

    关键步骤：
    1. 深拷贝XML元素
    2. 查找并复制所有关系引用（如图片blip）
    3. 将新关系ID更新到复制的XML中
    """
    try:
        # 深拷贝形状的XML元素
        new_elem = copy.deepcopy(shape._element)

        # 处理元素中的所有关系引用
        _remap_relationships(new_elem, source_part, target_part)

        # 添加到目标幻灯片
        target_spTree.append(new_elem)

    except Exception as e:
        logging.warning(f"复制形状失败: {e}")


def _remap_relationships(element, source_part, target_part):
    """
    重新映射元素中的所有关系引用

    处理的关系类型：
    - a:blip (图片)
    - a:hlinkClick (超链接)
    - r:link (外部链接)
    - c:chart (图表)
    等
    """
    # 定义需要处理的命名空间和属性
    EMBED_ATTRS = [
        (qn('r:embed'), RT.IMAGE),  # 图片嵌入
        (qn('r:link'), None),  # 外部链接
    ]

    LINK_ATTRS = [
        qn('r:id'),  # 通用关系ID
    ]

    # 收集所有需要重映射的rId
    rids_to_remap = set()

    # 查找所有带有 r:embed 属性的元素（主要是图片）
    for attr_name, rel_type in EMBED_ATTRS:
        for elem in element.iter():
            rid = elem.get(attr_name)
            if rid:
                rids_to_remap.add((elem, attr_name, rid))

    # 查找 a:blip 元素中的 r:embed（图片的主要引用方式）
    for blip in element.findall('.//' + qn('a:blip')):
        embed_rid = blip.get(qn('r:embed'))
        if embed_rid:
            rids_to_remap.add((blip, qn('r:embed'), embed_rid))
        link_rid = blip.get(qn('r:link'))
        if link_rid:
            rids_to_remap.add((blip, qn('r:link'), link_rid))

    # 查找超链接
    for hlinkClick in element.findall('.//' + qn('a:hlinkClick')):
        rid = hlinkClick.get(qn('r:id'))
        if rid:
            rids_to_remap.add((hlinkClick, qn('r:id'), rid))

    # 查找图表引用
    for chart in element.findall('.//' + qn('c:chart')):
        rid = chart.get(qn('r:id'))
        if rid:
            rids_to_remap.add((chart, qn('r:id'), rid))

    # 查找oleObject
    for oleObj in element.findall('.//' + qn('p:oleObj')):
        rid = oleObj.get(qn('r:id'))
        if rid:
            rids_to_remap.add((oleObj, qn('r:id'), rid))

    # 执行重映射
    for elem, attr_name, old_rid in rids_to_remap:
        try:
            new_rid = _copy_relationship(source_part, target_part, old_rid)
            if new_rid:
                elem.set(attr_name, new_rid)
        except Exception as e:
            logging.debug(f"重映射关系失败 {old_rid}: {e}")


def _copy_relationship(source_part, target_part, rid: str) -> Optional[str]:
    """
    复制关系及其目标资源

    返回新的关系ID
    """
    try:
        # 获取源关系
        source_rels = source_part.rels
        if rid not in source_rels:
            return rid  # 关系不存在，保持原样

        rel = source_rels[rid]

        # 检查是否是外部关系
        if rel.is_external:
            # 外部链接，直接创建新关系
            new_rid = target_part.relate_to(rel.target_ref, rel.reltype, is_external=True)
            return new_rid

        # 内部关系，需要复制目标资源
        target_resource = rel.target_part

        # 检查目标part中是否已经有相同的资源
        # 通过比较blob（二进制内容）来判断
        existing_rid = _find_existing_resource(target_part, target_resource)
        if existing_rid:
            return existing_rid

        # 复制资源并创建新关系
        new_rid = target_part.relate_to(target_resource, rel.reltype)
        return new_rid

    except Exception as e:
        logging.debug(f"复制关系 {rid} 失败: {e}")
        return rid  # 失败时保持原rId


def _find_existing_resource(target_part, resource_part) -> Optional[str]:
    """
    在目标part中查找是否已存在相同的资源

    通过比较资源的blob来判断
    """
    try:
        resource_blob = resource_part.blob

        for rid, rel in target_part.rels.items():
            if rel.is_external:
                continue
            try:
                if rel.target_part.blob == resource_blob:
                    return rid
            except:
                continue

        return None
    except:
        return None


def _copy_slide_background(source_slide, new_slide):
    """复制幻灯片背景"""
    try:
        source_cSld = source_slide._element.find(qn('p:cSld'))
        if source_cSld is not None:
            source_bg = source_cSld.find(qn('p:bg'))
            if source_bg is not None:
                new_cSld = new_slide._element.find(qn('p:cSld'))
                if new_cSld is not None:
                    # 移除已有背景
                    old_bg = new_cSld.find(qn('p:bg'))
                    if old_bg is not None:
                        new_cSld.remove(old_bg)

                    # 深拷贝背景
                    new_bg = copy.deepcopy(source_bg)

                    # 重映射背景中的关系（如背景图片）
                    _remap_relationships(new_bg, source_slide.part, new_slide.part)

                    # 插入到cSld的开头
                    new_cSld.insert(0, new_bg)
    except Exception as e:
        logging.debug(f"复制背景失败: {e}")


def _move_slide(prs: Presentation, from_idx: int, to_idx: int):
    """移动幻灯片位置"""
    try:
        sldIdLst = prs.slides._sldIdLst
        slides = list(sldIdLst)

        if from_idx < len(slides):
            slide = slides[from_idx]
            sldIdLst.remove(slide)

            # 调整目标索引
            actual_to = to_idx if to_idx < from_idx else to_idx - 1
            actual_to = min(actual_to, len(list(sldIdLst)))

            if actual_to >= len(list(sldIdLst)):
                sldIdLst.append(slide)
            else:
                sldIdLst.insert(actual_to, slide)
    except Exception as e:
        logging.warning(f"移动幻灯片失败: {e}")


def _build_shape_mapping_by_position(original_slide, new_slide) -> Dict[int, BaseShape]:
    """
    通过位置建立形状映射

    使用形状的位置和大小来匹配，比单纯按顺序更可靠
    """
    mapping = {}

    orig_shapes = list(original_slide.shapes)
    new_shapes = list(new_slide.shapes)

    # 为每个原始形状找到匹配的新形状
    used_new_indices = set()

    for orig_shape in orig_shapes:
        orig_geo = _get_shape_geometry(orig_shape)
        best_match = None
        best_distance = float('inf')
        best_idx = -1

        for idx, new_shape in enumerate(new_shapes):
            if idx in used_new_indices:
                continue

            new_geo = _get_shape_geometry(new_shape)

            # 计算位置和大小的差异
            distance = (
                    abs(orig_geo.left - new_geo.left) +
                    abs(orig_geo.top - new_geo.top) +
                    abs(orig_geo.width - new_geo.width) +
                    abs(orig_geo.height - new_geo.height)
            )

            # 如果完全匹配（或非常接近）
            if distance < best_distance:
                best_distance = distance
                best_match = new_shape
                best_idx = idx

        if best_match is not None and best_distance < Emu(Inches(0.1)):  # 容差0.1英寸
            mapping[orig_shape.shape_id] = best_match
            used_new_indices.add(best_idx)

            # 递归处理组合形状
            if orig_shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                try:
                    if best_match.shape_type == MSO_SHAPE_TYPE.GROUP:
                        _map_group_by_position(orig_shape, best_match, mapping)
                except:
                    pass

    # 对于未匹配的，尝试按顺序匹配
    unmatched_orig = [s for s in orig_shapes if s.shape_id not in mapping]
    unmatched_new = [s for i, s in enumerate(new_shapes) if i not in used_new_indices]

    for i, orig_shape in enumerate(unmatched_orig):
        if i < len(unmatched_new):
            mapping[orig_shape.shape_id] = unmatched_new[i]

    return mapping


def _map_group_by_position(orig_group, new_group, mapping: Dict[int, BaseShape]):
    """递归映射组合形状（基于位置）"""
    try:
        orig_sub = list(orig_group.shapes)
        new_sub = list(new_group.shapes)

        used_new_indices = set()

        for orig_shape in orig_sub:
            orig_geo = _get_shape_geometry(orig_shape)
            best_match = None
            best_distance = float('inf')
            best_idx = -1

            for idx, new_shape in enumerate(new_sub):
                if idx in used_new_indices:
                    continue

                new_geo = _get_shape_geometry(new_shape)
                distance = (
                        abs(orig_geo.left - new_geo.left) +
                        abs(orig_geo.top - new_geo.top)
                )

                if distance < best_distance:
                    best_distance = distance
                    best_match = new_shape
                    best_idx = idx

            if best_match is not None:
                mapping[orig_shape.shape_id] = best_match
                used_new_indices.add(best_idx)

                if orig_shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    if best_match.shape_type == MSO_SHAPE_TYPE.GROUP:
                        _map_group_by_position(orig_shape, best_match, mapping)

    except Exception as e:
        logging.debug(f"映射组合形状失败: {e}")


def _apply_blocks_with_mapping(slide, blocks: List[TextBlock],
                               shape_mapping: Dict[int, BaseShape],
                               target_lang: str,
                               slide_width: int, slide_height: int):
    """使用映射应用翻译"""

    # 按原始shape_id分组
    shape_blocks = _group_by_shape(blocks)

    # 获取几何信息
    all_geometries = _get_all_shape_geometries(slide)

    for orig_shape_id, sblocks in shape_blocks.items():
        if orig_shape_id not in shape_mapping:
            continue

        shape = shape_mapping[orig_shape_id]
        _apply_to_shape(shape, sblocks, target_lang, slide_width, slide_height, all_geometries)
